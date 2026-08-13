"""按用户提供的参考图重做风采展示 PPT(京东校招风格)

输出覆盖:京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示-终面版.pptx
"""
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = "/Users/fan.plain/Documents/工作/应聘简历"
SRC = f"{BASE}/京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示.pptx"
OUT = f"{BASE}/京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示-终面版.pptx"
AVATAR = f"{BASE}/作品集网站/assets/avatar.jpg"
QRCODE = f"{BASE}/作品集网站/assets/qrcode.png"

JD_RED = RGBColor(0xE1, 0x25, 0x1B)
JD_DEEP = RGBColor(0xC8, 0x16, 0x23)
INK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_RED = RGBColor(0xFD, 0xEF, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE8, 0xD8, 0xD7)

FONT = "Microsoft YaHei"
BAR_W = 1.15  # 左侧竖栏宽(英寸)


def style_run(run, size, bold=False, color=INK, font=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """lines: list of (text, size, bold, color, space_after)"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(lines):
        text, size, bold, color = item[:4]
        space_after = item[4] if len(item) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        style_run(run, size, bold, color)
    return box


def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=1.0, rounded=False, radius=0.12):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        rect.adjustments[0] = radius
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line_color is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line_color
        rect.line.width = Pt(line_w)
    rect.shadow.inherit = False
    return rect


def add_vertical_text(slide, x, y, w, h, lines, align=PP_ALIGN.CENTER):
    """竖排文本框(lines: (text, size, bold, color, space_after))"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("vert", "eaVert")
    for i, item in enumerate(lines):
        text, size, bold, color = item[:4]
        space_after = item[4] if len(item) > 4 else 6
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        style_run(run, size, bold, color)
    return box


def remove_shapes_with_text(slide, keywords):
    for sh in list(slide.shapes):
        if sh.has_text_frame and any(k in sh.text_frame.text for k in keywords):
            sh._element.getparent().remove(sh._element)


def side_bar(slide):
    """左侧京东校招竖栏"""
    add_rect(slide, 0, 0, BAR_W, 15.0, JD_DEEP)
    add_vertical_text(slide, 0.12, 0.45, 0.91, 9.6, [
        ("燃", 44, True, WHITE, 26),
        ("Jumpstart your dreams", 12, False, WHITE, 22),
        ("京东校招", 17, True, WHITE, 0),
    ])
    add_vertical_text(slide, 0.12, 10.6, 0.91, 4.0, [
        ("京东", 15, True, WHITE, 12),
        ("又好又便宜", 10.5, False, WHITE, 0),
    ])


def page_title(slide, title, subtitle=None):
    x = BAR_W + 0.55
    add_text(slide, x, 1.55, 24.6, 0.65, [(title, 23, True, JD_DEEP)])
    if subtitle:
        add_text(slide, x, 2.28, 24.6, 0.45, [(subtitle, 12, False, GRAY)])


prs = Presentation(SRC)
slides = list(prs.slides)

# ============================================================
# P1 封面
# ============================================================
s1 = slides[0]
side_bar(s1)
ph = None
for shp in s1.shapes:
    if shp.is_placeholder:
        ph = shp
        break
ph.left = Inches(BAR_W + 0.5)
ph.top = Inches(4.5)
ph.width = Inches(25.0 - BAR_W)
ph.height = Inches(1.6)
tf = ph.text_frame
tf.word_wrap = True
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "把复杂的事,做成确定的事"
style_run(r, 44, True, JD_RED)

add_text(s1, BAR_W + 0.5, 3.25, 25.0 - BAR_W, 0.5,
         [("京 东 TET · 综 合 方 向 · 个 人 风 采 展 示", 13, False, GRAY)], align=PP_ALIGN.CENTER)
add_text(s1, BAR_W + 0.5, 6.45, 25.0 - BAR_W, 0.75,
         [("赵奕帆", 27, True, INK)], align=PP_ALIGN.CENTER)
add_text(s1, BAR_W + 0.5, 7.35, 25.0 - BAR_W, 0.6,
         [("懂 AI、带过队、做成过生意 —— 一个想为京东供应链而生的复合型选手", 15, False, GRAY)],
         align=PP_ALIGN.CENTER)
s1.shapes.add_picture(AVATAR, Inches(12.0), Inches(8.45), Inches(2.9), Inches(2.9))
s1.shapes.add_picture(QRCODE, Inches(23.6), Inches(12.05), Inches(1.45), Inches(1.45))
add_text(s1, 22.95, 13.55, 2.75, 0.5, [("扫码查看我的作品集", 8.5, False, GRAY)], align=PP_ALIGN.CENTER)

# ============================================================
# P2 自我介绍
# ============================================================
s2 = slides[1]
remove_shapes_with_text(s2, ["个人自我介绍"])
side_bar(s2)
page_title(s2, "个人自我介绍")
nums = [
    ("6 年", "连续担任班长", "小学 · 初中 · 高中 · 本科 · 硕士"),
    ("1 个", "AI 智能问答产品", "LLM + RAG · 独立开发 · 三级级联架构"),
    ("15%", "采购成本压降", "6 家比价 · 建 3 家稳定渠道"),
]
for i, (num, label, note) in enumerate(nums):
    x = BAR_W + 1.2 + i * 6.1
    add_text(s2, x, 4.35, 5.7, 1.4, [(num, 62, True, JD_RED)], align=PP_ALIGN.CENTER)
    add_text(s2, x, 6.0, 5.7, 0.55, [(label, 16, True, INK)], align=PP_ALIGN.CENTER)
    add_text(s2, x, 6.7, 5.7, 0.55, [(note, 10.5, False, GRAY)], align=PP_ALIGN.CENTER)
add_text(s2, BAR_W + 1.2, 8.35, 17.2, 1.0,
         [("这三个数字,共同指向一句话:", 13.5, False, GRAY),
          ("把复杂的事,做成确定的事。", 17, True, JD_DEEP)])
s2.shapes.add_picture(AVATAR, Inches(20.6), Inches(4.1), Inches(4.3), Inches(4.3))

# ============================================================
# P3 校内外经历 / 项目经历(三卡片)
# ============================================================
s3 = slides[2]
remove_shapes_with_text(s3, ["校内外经历"])
side_bar(s3)
page_title(s3, "校内外经历 / 项目经历", "用高管面逻辑讲经历:背景 → 我的行动 → 结果 → 复盘迁移")

CX = BAR_W + 0.45
CARD_W = 7.55
GAP = 0.5
CARDS_Y = 3.1
CARDS_H = 8.5

p3_cards = [
    dict(
        num="01", title="创业项目负责人", sub="酶制剂创业方案 | 项目负责人 & 创始人",
        actions=[
            "判断产业机会,争取创业基金并组建 5 人团队",
            "制定 OKR 与周度经营复盘,统一研发/运营/市场节奏",
            "管理菌种、培养基、发酵辅料等上游采购",
        ],
        metric="-15%", metric_label="原料采购成本降低",
        metric_note="对 6 家供应商比价与实地考察,沉淀 3 家稳定渠道",
        migrate=["供应商开发与成本谈判", "价格谈判与经营复盘 · 经营管理与目标落地"],
    ),
    dict(
        num="02", title="AI 智能工单系统", sub="飞书CLI + OpenClaw + GPT/Gemini",
        actions=[
            "把多部门需求统一进入工单队列",
            "设计任务创建→分配→追踪→归档流程",
            "用优先级与资源占用可视化减少冲突",
        ],
        metric="+40%", metric_label="工单处理效率提升",
        metric_note="实验数据处理 / SOP检索 / 任务提醒自动化",
        migrate=["采购/运营协同提效", "流程标准化与自动化 · 跨部门协作与资源管理"],
    ),
    dict(
        num="03", title="功能注释 Agent", sub="多模态与大语言模型 | AI智能体开发",
        actions=[
            "独立完成数据采集、清洗、建模、可视化",
            "处理 9,170 个基因、12,353 个转录本",
            "三级级联推理架构 + Streamlit 看板",
        ],
        metric="21,523+", metric_label="结构化数据输出",
        metric_note="复杂数据管线可交互、可复用、可汇报",
        migrate=["经营数据看板与决策支持", "异常发现与快速追踪 · 数据驱动业务优化"],
    ),
]

for i, c in enumerate(p3_cards):
    x = CX + i * (CARD_W + GAP)
    add_rect(s3, x, CARDS_Y, CARD_W, CARDS_H, WHITE, line_color=JD_RED, line_w=1.4)
    # 编号块
    add_rect(s3, x + 0.3, CARDS_Y + 0.3, 0.55, 0.55, JD_RED, rounded=True, radius=0.3)
    add_text(s3, x + 0.3, CARDS_Y + 0.42, 0.55, 0.35, [(c["num"], 13, True, WHITE)], align=PP_ALIGN.CENTER)
    # 标题
    add_text(s3, x + 1.0, CARDS_Y + 0.3, CARD_W - 1.3, 0.6, [(c["title"], 17, True, INK)])
    add_text(s3, x + 1.0, CARDS_Y + 0.92, CARD_W - 1.3, 0.35, [(c["sub"], 9.5, False, GRAY)])
    y = CARDS_Y + 1.5
    # 关键动作
    add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.3, [("关键动作", 11.5, True, JD_RED)])
    y += 0.38
    for a in c["actions"]:
        add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.62, [("• " + a, 10.5, False, INK, 0)], line_spacing=1.05)
        y += 0.62
    # 分隔线
    add_rect(s3, x + 0.3, y + 0.04, CARD_W - 0.6, 0.016, LINE)
    y += 0.24
    # 验证结果
    add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.3, [("验证结果", 11.5, True, JD_RED)])
    y += 0.36
    add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.62, [(c["metric"], 26, True, JD_RED, 0)])
    y += 0.6
    add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.32, [(c["metric_label"], 11.5, True, INK)])
    y += 0.36
    add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.5, [(c["metric_note"], 9.5, False, GRAY)], line_spacing=1.0)
    y += 0.62
    add_rect(s3, x + 0.3, y + 0.04, CARD_W - 0.6, 0.016, LINE)
    y += 0.24
    # 京东迁移
    add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.3, [("京东迁移", 11.5, True, JD_RED)])
    y += 0.38
    for m in c["migrate"]:
        add_text(s3, x + 0.3, y, CARD_W - 0.6, 0.4, [(m, 9.8, False, INK, 0)], line_spacing=1.05)
        y += 0.42

# 底部横条
add_rect(s3, CX, 12.05, CARD_W * 3 + GAP * 2, 1.15, JD_DEEP, rounded=True, radius=0.5)
add_text(s3, CX, 12.05, CARD_W * 3 + GAP * 2, 1.15,
         [("我会这样讲经历:少讲「我学了很多」,多讲「约束 → 动作 → 结果 → 迁移」——结果可量化、可验证、可迁移,才是面试官关心的。",
           14, True, WHITE)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# P4 职业规划(三阶段 + 底部五项)
# ============================================================
s4 = slides[3]
remove_shapes_with_text(s4, ["职业规划"])
side_bar(s4)
page_title(s4, "职业规划", "不是空泛说「成为管理者」,而是到京东的经营场景:吃透一线 → 做结果 → 复制方法")

stages = [
    ("0-6 个月", "吃透一线", [
        "跑业务:选品、供应商、价盘、商家、履约",
        "跟数据:日/周/月经营数据,建立指标口径",
        "跟部门:看板、评价、优化闭环",
    ], "懂业务语言,拿出第一批改进点"),
    ("6-18 个月", "独立专项", [
        "承接一个跨范围项目:降本/提效/体验改善",
        "把 AI 工具用于数据清洗、问题定位、知识沉淀",
        "形成「问题—动作—结果」的复盘模板",
    ], "用结果证明方法,用方法复用结果"),
    ("18-36 个月", "带小团队", [
        "负责细分业务模块或区域的项目推进",
        "沉淀供应商策略/流程管理/经营分析方法论",
        "从执行者成长为带教人,带项目,带组织能力",
    ], "带团队拿结果,输出可复制能力"),
]
SW = 7.35
SGAP = 0.52
SY = 3.2
SH = 5.9
add_text(s4, CX, SY - 0.42, 10, 0.4, [("三阶段成长路径", 13, True, JD_RED)])
for i, (phase, name, items, goal) in enumerate(stages):
    x = CX + i * (SW + SGAP)
    add_rect(s4, x, SY, SW, SH, WHITE, line_color=JD_RED, line_w=1.4)
    add_rect(s4, x, SY, SW, 0.78, JD_RED)
    add_text(s4, x + 0.28, SY + 0.12, SW - 0.56, 0.55,
             [(f"{phase} | {name}", 15, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s4, x + 0.28, SY + 1.0, SW - 0.56, 3.2,
             [(t, 10.8, False, INK, 8) for t in items], line_spacing=1.08)
    add_rect(s4, x + 0.28, SY + SH - 1.15, SW - 0.56, 0.9, LIGHT_RED, rounded=True, radius=0.14)
    add_text(s4, x + 0.4, SY + SH - 1.15, SW - 0.8, 0.9,
             [("目标:" + goal, 10.5, True, JD_DEEP, 0)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    if i < 2:
        ax = x + SW + 0.06
        add_text(s4, ax, SY + SH / 2 - 0.3, 0.4, 0.6, [("→", 24, True, JD_RED)], align=PP_ALIGN.CENTER)

add_text(s4, CX, 9.7, 24.4, 0.45,
         [("我理解京东业务的底层不是「喊口号」,而是围绕成本、效率、产品、价格、服务做连续改善", 13, True, INK)],
         align=PP_ALIGN.CENTER)

pillars = [
    ("成本", "供应链谈判 / 采购降本"),
    ("效率", "工具提效 / AI自动化"),
    ("产品", "数据洞察 / 质量优化"),
    ("价格", "比价+动态定价策略"),
    ("服务", "从用户满意倒推链路动作"),
]
PW = 4.6
PX0 = CX
for i, (name, desc) in enumerate(pillars):
    x = PX0 + i * (PW + 0.16)
    add_rect(s4, x, 10.55, PW, 1.75, LIGHT_RED, line_color=JD_RED, line_w=1.0, rounded=True, radius=0.1)
    add_text(s4, x + 0.2, 10.85, PW - 0.4, 0.45, [(name, 14, True, JD_RED)], align=PP_ALIGN.CENTER)
    add_text(s4, x + 0.2, 11.42, PW - 0.4, 0.7, [(desc, 10.5, False, INK)], align=PP_ALIGN.CENTER, line_spacing=1.0)

# ============================================================
# P5 为什么选择加入京东
# ============================================================
s5 = slides[4]
remove_shapes_with_text(s5, ["为什么选择加入京东"])
side_bar(s5)
page_title(s5, "为什么选择加入京东", "我的选择逻辑:京东的供应链业务土壤,能把实践的经验、技术和组织能力真正放大")

# 左上框
LX = CX
LW = 12.15
LH = 4.15
add_rect(s5, LX, 3.0, LW, LH, WHITE, line_color=JD_RED, line_w=1.4)
add_text(s5, LX + 0.35, 3.25, LW - 0.7, 0.4, [("我对京东和 TET 的理解", 14.5, True, JD_DEEP)])
add_text(s5, LX + 0.35, 3.8, LW - 0.7, 3.2, [
    ("京东定位:以供应链为基础的技术与服务企业,长期主义,客户为先", 11, False, INK, 7),
    ("TET 目标:培养真正具备业务成事、推动战略落地的中高层管理者", 11, False, INK, 7),
    ("综合方向:深入一线、多场景、复杂运营;快速学习、跨部门协同拿结果", 11, False, INK, 0),
], line_spacing=1.1)
# 右上框
RX = LX + LW + 0.65
add_rect(s5, RX, 3.0, LW, LH, WHITE, line_color=JD_RED, line_w=1.4)
add_text(s5, RX + 0.35, 3.25, LW - 0.7, 0.65,
         [("所以选择京东,不是因为「平台大」", 14.5, True, JD_DEEP)],
         line_spacing=1.0)
add_text(s5, RX + 0.35, 3.9, LW - 0.7, 3.1, [
    ("而是因为这里的问题足够真实、链路足够长、结果足够可衡量", 10.5, False, GRAY, 8),
    ("真实业务场景:接触供应商、用户、商品、履约全链路", 11, False, INK, 6),
    ("价值能证明:数据说话,结果可衡量,贡献可复盘", 11, False, INK, 6),
    ("成长可复制:沉下去做透,再把方法复制到更大范围", 11, False, INK, 0),
], line_spacing=1.1)

# 中下框:价值观
VY = 7.55
VH = 4.1
add_rect(s5, LX, VY, LW * 2 + 0.65, VH, WHITE, line_color=JD_RED, line_w=1.4)
add_text(s5, LX + 0.35, VY + 0.22, LW * 2 - 0.05, 0.4,
         [("价值观匹配:用经历证明,而不是口头表态", 14.5, True, JD_DEEP)])
values = [
    ("客户为先", "从用户业务痛点倒推方案"),
    ("创新", "AI Agent/工具/数据看板落地提效"),
    ("拼搏担当", "0→1 创业项目负责人,目标导向、敢打硬仗"),
    ("诚信", "供应商比价、长期合作与规则意识"),
    ("感恩", "被师兄托举过,愿意从一线学起、把善意传下去"),
]
VW = 4.6
for i, (name, desc) in enumerate(values):
    x = LX + 0.35 + i * (VW + 0.08)
    add_rect(s5, x, VY + 0.8, VW, VH - 1.35, LIGHT_RED, rounded=True, radius=0.08)
    add_text(s5, x + 0.2, VY + 1.05, VW - 0.4, 0.45, [(name, 13.5, True, JD_RED)], align=PP_ALIGN.CENTER)
    add_text(s5, x + 0.2, VY + 1.62, VW - 0.4, 1.6, [(desc, 10.2, False, INK)], align=PP_ALIGN.CENTER, line_spacing=1.08)

# 底部横条
BY = 12.05
add_rect(s5, LX, BY, LW * 2 + 0.65, 1.1, JD_DEEP, rounded=True, radius=0.5)
add_text(s5, LX, BY, LW * 2 + 0.65, 1.1,
         [("入职承诺:不空谈技术、不空谈管理;先把一线业务做扎实,再用方法和工具放大结果。", 14, True, WHITE)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("saved:", OUT)
