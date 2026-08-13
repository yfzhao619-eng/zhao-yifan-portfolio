"""检查风采展示 PPT 模板结构"""
from pptx import Presentation
from pptx.util import Emu

SRC = "/Users/fan.plain/Documents/工作/应聘简历/京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示.pptx"

prs = Presentation(SRC)
print("slide size:", Emu(prs.slide_width).inches, "x", Emu(prs.slide_height).inches, "inches")
print("layouts:", [l.name for l in prs.slide_masters[0].slide_layouts])
for i, slide in enumerate(prs.slides):
    print(f"===== Slide {i+1} layout={slide.slide_layout.name} =====")
    for j, shape in enumerate(slide.shapes):
        print(
            f"  [{j}] id={shape.shape_id} type={shape.shape_type} name={shape.name!r} "
            f"pos=({Emu(shape.left).inches:.2f},{Emu(shape.top).inches:.2f}) "
            f"size=({Emu(shape.width).inches:.2f}x{Emu(shape.height).inches:.2f})"
        )
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs)
                if t.strip():
                    print(f"       text: {t!r}")
        elif shape.shape_type == 13:
            print("       [picture]")
