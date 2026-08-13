"""校验 PPT 所有形状是否越界,并输出每页元素数量"""
from pptx import Presentation
from pptx.util import Emu

SRC = "/Users/fan.plain/Documents/工作/应聘简历/京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示-终面版.pptx"
prs = Presentation(SRC)
PW = prs.slide_width
PH = prs.slide_height
issues = []
for i, slide in enumerate(prs.slides):
    n = len(slide.shapes)
    texts = 0
    for sh in slide.shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            continue
        if l is not None and (l < -Emu(91440) or t < -Emu(91440) or l + w > PW + Emu(91440) or t + h > PH + Emu(91440)):
            issues.append(
                f"slide {i+1}: {sh.shape_type} out of bounds "
                f"({Emu(l).inches:.2f},{Emu(t).inches:.2f},{Emu(w).inches:.2f},{Emu(h).inches:.2f})"
            )
        if sh.has_text_frame:
            texts += 1
    print(f"slide {i+1}: {n} shapes, {texts} with text")
print("---")
if issues:
    print("\n".join(issues))
else:
    print("no out-of-bounds shapes")
