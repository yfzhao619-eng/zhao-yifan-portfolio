"""基于现有风采展示模板生成终面版 PPT(不修改原文件)

输出:京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示-终面版.pptx
"""
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = "/Users/fan.plain/Documents/工作/应聘简历"
SRC = f"{BASE}/京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示.pptx"
OUT = f"{BASE}/京东/TET/TET面试材料汇总/TET21+综合方向+赵奕帆+风采展示-终面版.pptx"
AVATAR = f"{BASE}/作品集网站/assets/avatar.jpg"
QRCODE = f"{BASE}/作品集网站/assets/qrcode.png"

JD_RED = RGBColor(0xE1, 0x25, 0x1B)
JD_DEEP = RGBColor(0xC8, 0x16, 0x23)
INK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_RED = RGBColor(0xFD, 0xEF, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xD9, 0xA4, 0x41)

FONT = "Microsoft YaHei"


def style_run(run, size, bold=False, color=INK, font=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) 或 (text, size, bold, color, space_after)"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(lines):
        text, size, bold, color = item[:4]
        space_after = item[4] if len(item) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        style_run(run, size, bold, color)
    return box


def add_card(slide, x, y, w, h, fill=WHITE, line_color=None, radius=0.055):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.adjustments[0] = radius
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line_color is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = line_color
        card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


prs = Presentation(SRC)
slides = list(prs.slides)

# ================= P1 封面 =================
s1 = slides[0]
ph = None
for shp in s1.shapes:
    if shp.is_placeholder:
        ph = shp
        break
assert ph is not None, "P1 占位符未找到"
ph.left = Inches(2.0)
ph.top = Inches(4.7)
ph.width = Inches(22.7)
ph.height = Inches(1.7)
tf = ph.text_frame
tf.word_wrap = True
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "把复杂的事,做成确定的事"
style_run(r, 46, True, JD_RED)

add_text(s1, 2.0, 3.3, 22.7, 0.6,
         [("京 东 TET · 综 合 方 向 · 个 人 风 采 展 示", 14, False, GRAY)],
         align=PP_ALIGN.CENTER)
add_text(s1, 2.0, 6.7, 22.7, 0.8,
         [("赵奕帆", 28, True, INK)],
         align=PP_ALIGN.CENTER)
add_text(s1, 2.0, 7.7, 22.7, 0.7,
         [("懂 AI、带过队、做成过生意 —— 一个想为京东供应链而生的复合型选手", 16, False, GRAY)],
         align=PP_ALIGN.CENTER)
# 头像
s1.shapes.add_picture(AVATAR, Inches(11.84), Inches(8.9), Inches(3.0), Inches(3.0))
# 二维码
s1.shapes.add_picture(QRCODE, Inches(23.3), Inches(11.9), Inches(1.5), Inches(1.5))
add_text(s1, 22.7, 13.45, 2.7, 0.7,
         [("扫码查看我的作品集", 9, False, GRAY)],
         align=PP_ALIGN.CENTER)

# ================= P2 自我介绍 =================
s2 = slides[1]
# 三个数字
nums = [
    ("6 年", "连续担任班长", "小学 · 初中 · 高中 · 本科 · 硕士"),
    ("1 个", "AI 智能问答产品", "LLM + RAG · 独立开发 · 三级级联架构"),
    ("15%", "采购成本压降", "6 家比价 · 建 3 家稳定渠道"),
]
col_x = [1.6, 7.7, 13.8]
for i, (num, label, note) in enumerate(nums):
    x = col_x[i]
    add_text(s2, x, 4.5, 5.6, 1.5, [(num, 66, True, JD_RED)], align=PP_ALIGN.CENTER)
    add_text(s2, x, 6.15, 5.6, 0.6, [(label, 17, True, INK)], align=PP_ALIGN.CENTER)
    add_text(s2, x, 6.9, 5.6, 0.6, [(note, 11, False, GRAY)], align=PP_ALIGN.CENTER)
# 底部一句话
add_text(s2, 1.6, 8.6, 16.4, 1.0,
         [("这三个数字,共同指向一句话:", 14, False, GRAY),
          ("把复杂的事,做成确定的事。", 18, True, JD_DEEP)])
# 证件照
s2.shapes.add_picture(AVATAR, Inches(20.2), Inches(4.2), Inches(4.4), Inches(4.4))

# ================= P3 经历故事链 =================
s3 = slides[2]
# 更新标题
title_shape = None
for sh in s3.shapes:
    if sh.has_text_frame and "校内外经历" in sh.text_frame.text:
        title_shape = sh
if title_shape is not None:
    title_shape.text_frame.clear()
    p = title_shape.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "我的经历如何塑造我"
    style_run(r, 24, True, JD_DEEP)
add_text(s3, 1.6, 3.2, 23.5, 0.5,
         [("一条故事链:把复杂的事,做成确定的事", 14, False, GRAY)])
cards = [
    ("① 管理者底色", "6 年班长 × 组织部", "200+ 人团日 · 500+ 人音乐节 · 经费 -15%"),
    ("② 技术硬实力", "AI Agent 开发", "9,170 基因 / 12,353 转录本 · 提速数万倍"),
    ("③ 商业实践", "酶制剂创业 · 项目负责人", "6 家比价 → 3 家渠道 · 采购成本 -15%"),
    ("④ 组织提效", "AI 智能工单系统", "效率 +40% · 零冲突零遗漏"),
]
cw, ch = 11.55, 3.6
xs = [1.6, 13.55]
ys = [3.95, 7.95]
for i, (tag, title, metrics) in enumerate(cards):
    x = xs[i % 2]
    y = ys[i // 2]
    add_card(s3, x, y, cw, ch, fill=LIGHT_RED)
    add_text(s3, x + 0.5, y + 0.35, cw - 1.0, 0.5,
             [(tag, 15, True, JD_RED)])
    add_text(s3, x + 0.5, y + 1.0, cw - 1.0, 0.6,
             [(title, 20, True, INK)])
    add_text(s3, x + 0.5, y + 1.95, cw - 1.0, 1.2,
             [(metrics, 15, False, GRAY)])
# 底部故事链主线
add_card(s3, 1.6, 11.95, 23.5, 1.1, fill=JD_DEEP)
add_text(s3, 1.6, 11.95, 23.5, 1.1,
         [("管理拆解复杂  →  AI 压缩复杂  →  供应链交付确定  →  系统重塑协作", 18, True, WHITE)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================= P4 职业规划 =================
s4 = slides[3]
add_card(s4, 1.6, 3.3, 11.8, 9.3, fill=WHITE, line_color=JD_RED)
add_text(s4, 2.0, 3.7, 11.0, 0.7,
         [("业务方向:京东零售采销 × AI 提效", 18, True, JD_DEEP)])
add_text(s4, 2.0, 4.7, 11.0, 7.6, [
    ("▸ 战略认知:2025 年京东拆分采销单元,资源压到离商品和用户最近的地方", 13.5, False, INK, 10),
    ("▸ 确定性心智 = 采销 + 物流共同交付,采销是第一道闸门", 13.5, False, INK, 10),
    ("▸ AI 落地价值最大的场景:选品 / 比价 / 定价 / 补货,决策高频、数据丰富、结果可量化", 13.5, False, INK, 10),
    ("▸ 我最想做的事:建设「AI 采销」——智能寻源比价、需求预测、多供应商多品类并行管理", 13.5, True, JD_RED, 10),
])
add_card(s4, 14.0, 3.3, 11.1, 9.3, fill=WHITE, line_color=JD_RED)
add_text(s4, 14.4, 3.7, 10.3, 0.7,
         [("管理路径:为做成这件事", 18, True, JD_DEEP)])
add_text(s4, 14.4, 4.7, 10.3, 7.6, [
    ("轮岗期(1.5~2 年)· 一线业务学生 + AI 提效先行者", 14, True, INK, 6),
    ("   补:品类知识、采销全链路、成本结构、履约仓配逻辑", 12.5, False, GRAY, 8),
    ("   做:把重复劳动工具化,积累 2~3 个 AI 提效案例", 12.5, False, GRAY, 14),
    ("定岗后(2~3 年)· 业务 Owner + AI 项目发起人", 14, True, INK, 6),
    ("   补:品类损益、谈判进阶、供应商生态、跨部门协调", 12.5, False, GRAY, 8),
    ("   做:扛一个品类经营,发起 AI 采销专项", 12.5, False, GRAY, 14),
    ("长期:既懂采销业务、又能用 AI 重构业务的经营管理者", 13.5, True, JD_RED, 0),
])

# ================= P5 为什么京东 =================
s5 = slides[4]
cols = [
    ("京东给我的平台", "供应链",
     "创业时只有 6 家供应商的「微缩版」;京东拥有 1500+ 仓、一体化供应链的「完整版」——同一个方法论,千倍的放大空间"),
    ("我依托京东的", "三样东西",
     "① 供应链基础设施:确定性的物理底座\n② 自营品质心智:正品 · 售后无忧\n③ AI 战略投入:言犀 · JoyAI"),
    ("为什么是 TET", "训练场",
     "我要的不是一个岗位,而是「从业务一线到经营全局」的训练场;轮岗 + 高管带教,补上大组织全局视野"),
]
cw2 = 7.5
for i, (h1, h2, body) in enumerate(cols):
    x = 1.6 + i * (cw2 + 0.55)
    add_card(s5, x, 3.2, cw2, 6.9, fill=WHITE, line_color=JD_RED)
    add_text(s5, x + 0.4, 3.55, cw2 - 0.8, 0.6, [(h1, 18, True, JD_DEEP)])
    add_text(s5, x + 0.4, 4.4, cw2 - 0.8, 0.5, [(h2, 13, False, GRAY)])
    add_text(s5, x + 0.4, 5.1, cw2 - 0.8, 4.8, [(body, 13, False, INK, 8)])
# 价值观条
add_card(s5, 1.6, 10.5, 23.5, 1.35, fill=LIGHT_RED)
add_text(s5, 1.6, 10.5, 23.5, 1.35, [
    ("客户为先 · 创新 · 拼搏 · 担当 · 感恩 · 诚信", 20, True, JD_DEEP),
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s5, 1.6, 12.15, 23.5, 0.8,
         [("感恩不是一句口号,是一种循环 —— 被师兄托举过的我,会本能地去托举别人", 14, False, GRAY)],
         align=PP_ALIGN.CENTER)

prs.save(OUT)
print("saved:", OUT)
